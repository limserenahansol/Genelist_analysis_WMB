#!/usr/bin/env python
"""Static layout check for the planning deck.

No renderer is installed on this machine, so overflow is detected geometrically
instead: tables auto-grow their rows in PowerPoint, which is the failure mode
that silently pushes content off the bottom of a slide. Estimated text heights
are approximate, so treat the output as a list of places to eyeball rather than
as ground truth.
"""
from __future__ import annotations

import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

PPTX = Path(__file__).resolve().parents[2] / "outputs" / "NextPhase_Plan_ORBm_BMAp_TRAP.pptx"

# Calibri averages a bit over half an em per character across mixed-case text.
CHAR_EM = 0.50
LINE_FACTOR = 1.22
CELL_MARGIN_X = Inches(0.1) * 2
CELL_MARGIN_Y = Inches(0.05) * 2


def text_lines(text: str, width_emu: int, pt: float) -> int:
    usable = max(Emu(1), width_emu - CELL_MARGIN_X)
    char_w = Pt(pt * CHAR_EM)
    per_line = max(1, int(usable / char_w))
    total = 0
    for seg in (text or "").split("\n"):
        total += max(1, math.ceil(len(seg) / per_line))
    return total


def cell_font_pt(cell, default=10.0) -> float:
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                return r.font.size.pt
    return default


def check(prs) -> list[str]:
    problems: list[str] = []
    sh, sw = prs.slide_height, prs.slide_width

    for idx, slide in enumerate(prs.slides, start=1):
        boxes = []
        for shape in slide.shapes:
            top, left = shape.top or 0, shape.left or 0
            width = shape.width or 0
            height = shape.height or 0

            if shape.has_table:
                tbl = shape.table
                col_w = [c.width for c in tbl.columns]
                needed = 0
                for ri, row in enumerate(tbl.rows):
                    tallest = row.height
                    for ci, cell in enumerate(row.cells):
                        pt = cell_font_pt(cell)
                        n = text_lines(cell.text, col_w[ci], pt)
                        h = int(Pt(n * pt * LINE_FACTOR)) + CELL_MARGIN_Y
                        tallest = max(tallest, h)
                    needed += tallest
                    if tallest > row.height * 1.35 and tallest - row.height > Inches(0.08):
                        problems.append(
                            f"  slide {idx}: table row {ri} wants "
                            f"{tallest / 914400:.2f}\" but is set to {row.height / 914400:.2f}\""
                        )
                if needed > height:
                    problems.append(
                        f"  slide {idx}: TABLE GROWS {(needed - height) / 914400:.2f}\" "
                        f"(declared {height / 914400:.2f}\" -> ~{needed / 914400:.2f}\")"
                    )
                height = max(height, needed)

            bottom, right = top + height, left + width
            if bottom > sh + Inches(0.02):
                problems.append(
                    f"  slide {idx}: {shape.shape_type} '{(shape.name or '')[:22]}' "
                    f"overflows BOTTOM by {(bottom - sh) / 914400:.2f}\""
                )
            if right > sw + Inches(0.02):
                problems.append(
                    f"  slide {idx}: {shape.shape_type} '{(shape.name or '')[:22]}' "
                    f"overflows RIGHT by {(right - sw) / 914400:.2f}\""
                )
            if height and width:
                boxes.append((shape.name or "?", left, top, right, bottom, bool(shape.has_table)))

        # Overlap between content blocks. Background bands and the title rule are
        # full-width by design, so only compare blocks that both sit inside the body.
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                if min(a[3], b[3]) - max(a[1], b[1]) <= Inches(0.05):
                    continue
                ov = min(a[4], b[4]) - max(a[2], b[2])
                if ov <= Inches(0.05):
                    continue
                if (a[3] - a[1]) > Inches(12) or (b[3] - b[1]) > Inches(12):
                    continue
                if not (a[5] or b[5]):
                    continue
                problems.append(
                    f"  slide {idx}: '{a[0][:20]}' and '{b[0][:20]}' overlap "
                    f"vertically by {ov / 914400:.2f}\""
                )
    return problems


def main() -> int:
    prs = Presentation(PPTX)
    print(f"{PPTX.name}: {len(prs.slides._sldIdLst)} slides, "
          f"{prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} in\n")
    problems = check(prs)
    if not problems:
        print("No geometry problems detected.")
        return 0
    print(f"{len(problems)} thing(s) to look at:")
    for p in problems:
        print(p)
    return 1


if __name__ == "__main__":
    sys.exit(main())
