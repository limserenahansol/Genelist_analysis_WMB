#!/usr/bin/env python
"""
make_next_phase_deck.py

Builds the PI-facing deck proposing the next phase of the opioid TRAP project:
Experiment 1 (Xenium on the Post-tagged ensemble in ORBm + BMAp) and
Experiment 2 (TRAP2 chemogenetic inhibition of that ensemble).

Every timing number in the deck traces to a citation in REFERENCES; the
scheduling logic is derived in slides 8-9 rather than asserted.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1F, 0x3B, 0x63)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xEE, 0xF3, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB3, 0x1B, 0x1B)
GREEN = RGBColor(0x1E, 0x6F, 0x3C)
AMBER = RGBColor(0x9C, 0x63, 0x00)
BAND = RGBColor(0xD9, 0xE4, 0xF0)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.55)
BODY_TOP = Inches(1.42)
BODY_H = Inches(5.45)

# Usable content band. Right-hand columns are derived from RIGHT_EDGE so nothing
# runs off the slide; laying them out by eye put several panels ~0.5" over.
RIGHT_EDGE = W - M
CONTENT_W = RIGHT_EDGE - M
COL2_W = Inches(5.96)
COL2_L = RIGHT_EDGE - COL2_W
COL3_W = Inches(3.88)


def right_of(table_w, gap=Inches(0.3)):
    """Left edge and width of the panel column beside a table of width table_w."""
    left = M + table_w + gap
    return left, RIGHT_EDGE - left


def deck() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _tf(shape, size, color=None, bold=False, space_after=4, line=0.95):
    tf = shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.space_after = Pt(space_after)
        para.line_spacing = line
        for r in para.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            if color:
                r.font.color.rgb = color
    return tf


def title(slide, text, sub=None, kicker=None):
    if kicker:
        k = slide.shapes.add_textbox(M, Inches(0.22), W - 2 * M, Inches(0.28))
        k.text_frame.text = kicker.upper()
        _tf(k, 11, BLUE, bold=True)
    t = slide.shapes.add_textbox(M, Inches(0.5), W - 2 * M, Inches(0.55))
    t.text_frame.text = text
    _tf(t, 25, NAVY, bold=True)
    if sub:
        s = slide.shapes.add_textbox(M, Inches(1.03), W - 2 * M, Inches(0.36))
        s.text_frame.text = sub
        _tf(s, 13, GREY)
    ln = slide.shapes.add_shape(1, M, Inches(1.35), W - 2 * M, Emu(11430))
    ln.fill.solid()
    ln.fill.fore_color.rgb = BAND
    ln.line.fill.background()
    ln.shadow.inherit = False


def bullets(slide, items, left, top, width, height, size=13, gap=7):
    """items: (text, level, color, bold)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        text, lvl, col, bold = (list(it) + [0, None, False])[:4]
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = lvl
        para.space_after = Pt(gap if lvl == 0 else max(2, gap - 3))
        para.line_spacing = 0.98
        r = para.add_run()
        r.text = ("" if lvl == 0 else "– ") + text
        r.font.size = Pt(size if lvl == 0 else size - 1.5)
        r.font.bold = bold
        r.font.color.rgb = col or (NAVY if bold else GREY)
    return box


def table(slide, rows, left, top, width, col_w=None, size=11, hdr_size=11,
          hdr_fill=NAVY, row_h=Inches(0.3), hdr_h=Inches(0.34), align=None,
          cell_colors=None):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, left, top, width, hdr_h + row_h * (nr - 1)).table
    if col_w:
        total = sum(col_w)
        for j, w in enumerate(col_w):
            gt.columns[j].width = Emu(int(width * w / total))
    gt.rows[0].height = hdr_h
    for i in range(1, nr):
        gt.rows[i].height = row_h
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.text = str(val)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.06)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = hdr_fill
            else:
                c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            if cell_colors and (i, j) in cell_colors:
                c.fill.fore_color.rgb = cell_colors[(i, j)]
            for para in c.text_frame.paragraphs:
                para.line_spacing = 0.92
                if align and j < len(align):
                    para.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                                      "r": PP_ALIGN.RIGHT}[align[j]]
                for r in para.runs:
                    r.font.size = Pt(hdr_size if i == 0 else size)
                    r.font.bold = i == 0
                    r.font.color.rgb = WHITE if i == 0 else NAVY
    return gt


def panel(slide, left, top, width, height, head, body, accent=BLUE, size=12):
    bar = slide.shapes.add_shape(1, left, top, width, Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.text = head
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    _tf(bar, 11.5, WHITE, bold=True)
    bg = slide.shapes.add_shape(1, left, top + Inches(0.32), width, height - Inches(0.32))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.color.rgb = BAND
    bg.shadow.inherit = False
    tf = bg.text_frame
    tf.word_wrap = True
    # Autoshapes default to centred text, which silently centres the first
    # paragraph of every panel; anchor and align explicitly.
    bg.vertical_anchor = MSO_ANCHOR.TOP
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.09)
    for i, (txt, bold) in enumerate(body):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(5)
        para.line_spacing = 0.97
        r = para.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = NAVY if bold else GREY
    return bg


def callout(slide, left, top, width, height, text, fill=RGBColor(0xFF, 0xF4, 0xD6),
            edge=AMBER, size=12.5, bold=True):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = edge
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.06)
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.LEFT
    _tf(box, size, NAVY, bold=bold, line=1.0)
    return box


MORPH = RGBColor(0xC0, 0x50, 0x4D)
WATER = RGBColor(0x7A, 0xA6, 0xC2)
NEUTRAL = RGBColor(0xB7, 0xB7, 0xB7)
ABST = RGBColor(0xD9, 0xD2, 0xE9)


def phase_timeline(slide, left, top, width, segments, height=Inches(0.42),
                   marks=None, label_size=9.5, day_size=8.5):
    """segments: (label, n_days, fill). marks: (day_center_fraction, text, colour)."""
    total = sum(sg[1] for sg in segments)
    x = left
    for label, days, fill in segments:
        w = Emu(int(width * days / total))
        box = slide.shapes.add_shape(1, x, top, w, height)
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = WHITE
        box.line.width = Pt(1.25)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.02)
        tf.margin_top = tf.margin_bottom = 0
        tf.text = label
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.line_spacing = 0.85
            for r in para.runs:
                r.font.size = Pt(label_size)
                r.font.bold = True
                r.font.color.rgb = WHITE if fill in (MORPH, NAVY, GREEN) else NAVY
        x += w
    for frac, text, col in (marks or []):
        cx = left + Emu(int(width * frac))
        pin = slide.shapes.add_shape(1, cx - Emu(9000), top - Inches(0.14),
                                     Emu(18000), Inches(0.14))
        pin.fill.solid()
        pin.fill.fore_color.rgb = col
        pin.line.fill.background()
        pin.shadow.inherit = False
        lab = slide.shapes.add_textbox(cx - Inches(0.85), top - Inches(0.52),
                                       Inches(1.7), Inches(0.36))
        lab.text_frame.word_wrap = True
        lab.text_frame.text = text
        for para in lab.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.line_spacing = 0.85
            for r in para.runs:
                r.font.size = Pt(day_size)
                r.font.bold = True
                r.font.color.rgb = col


def legend(slide, left, top, items, size=9):
    x = left
    for label, col in items:
        sw = slide.shapes.add_shape(1, x, top + Inches(0.03), Inches(0.16), Inches(0.13))
        sw.fill.solid()
        sw.fill.fore_color.rgb = col
        sw.line.fill.background()
        sw.shadow.inherit = False
        tb = slide.shapes.add_textbox(x + Inches(0.2), top - Inches(0.02),
                                      Inches(0.12 * len(label) + 0.4), Inches(0.24))
        tb.text_frame.text = label
        _tf(tb, size, GREY)
        x += Inches(0.2) + Inches(0.075 * len(label) + 0.18)


def footer(slide, text):
    f = slide.shapes.add_textbox(M, H - Inches(0.42), W - 2 * M, Inches(0.28))
    f.text_frame.text = text
    _tf(f, 9.5, GREY)
    f.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


# ----------------------------------------------------------------------------


def s01_title(prs):
    s = blank(prs)
    bg = s.shapes.add_shape(1, 0, 0, W, Inches(2.9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False

    t = s.shapes.add_textbox(M, Inches(0.85), W - 2 * M, Inches(1.0))
    t.text_frame.text = ("The Post-active ORBm and BMAp ensemble:\n"
                         "what it is, and whether relapse needs it")
    _tf(t, 30, WHITE, bold=True, line=1.05)

    st = s.shapes.add_textbox(M, Inches(2.05), W - 2 * M, Inches(0.5))
    st.text_frame.text = ("Proposal for the next phase — two experiments on the two regions "
                          "our whole-brain TRAP screen selected")
    _tf(st, 14, RGBColor(0xC9, 0xDA, 0xEC))

    meta = s.shapes.add_textbox(M, Inches(3.25), Inches(5.6), Inches(1.2))
    meta.text_frame.text = "Hansol Lim\nMorphine lick-PR / TRAP2;Ai14 / Xenium\nAugust 2026"
    _tf(meta, 13, NAVY, bold=False, line=1.3)

    panel(s, COL2_L, Inches(3.15), COL2_W, Inches(3.3),
          "WHAT I AM ASKING FOR",
          [("1.  Approval to run the Xenium pilot with Alina's lab in September, "
            "using the base mouse brain panel on tissue we already have.", False),
           ("2.  Approval to order the panel in October: base Mouse Brain panel plus an add-on. "
            "I have matched our gene list against the real base panel, and it needs only "
            "51 of the 100 custom slots, so this is the cheaper of the two routes.", False),
           ("3.  Agreement on one design change to the chemogenetics plan: "
            "test at protracted abstinence rather than inside the 3-day Post window. "
            "Slides 8-9 show why the original plan would have produced a false negative.", False),
           ("4.  Sign-off on group sizes. My first sketch (n=4) has ~29% power; "
            "I propose it as an expression pilot and a powered cohort behind it.", False)],
          accent=BLUE, size=12)
    footer(s, "All timing parameters cited on the final slide.")


def s02_where_we_are(prs):
    s = blank(prs)
    title(s, "Where we are", "Three pieces of work are finished and feed directly into the next phase",
          kicker="Status")

    panel(s, M, BODY_TOP, COL3_W, Inches(3.1), "1  WHOLE-BRAIN TRAP SCREEN  (done)",
          [("20 mice, 10 Active / 10 Passive, terminal design: each mouse is "
            "sacrificed at one behavioural phase.", False),
           ("Phases sampled: During, Post, Withdrawal, Reinstatement.", False),
           ("Readout: tdTomato+ cell density per Allen region, "
            "Active vs Passive within phase.", False)], accent=BLUE)

    panel(s, M + COL3_W + Inches(0.29), BODY_TOP, COL3_W, Inches(3.1),
          "2  REGION SELECTION  (done)",
          [("7 candidate regions carried forward for spatial transcriptomics: "
            "ORBm, AId, CA, BMAp, LM, RE, CP.", False),
           ("Selection rule: Active > Passive in the craving phases "
            "(Post and Reinstatement) and still separated during Withdrawal.", False),
           ("ORBm and BMAp are the two I propose to take forward now.", True)], accent=BLUE)

    panel(s, M + 2 * (COL3_W + Inches(0.29)), BODY_TOP, COL3_W, Inches(3.1),
          "3  PROBE PANEL  (done)",
          [("Built from Allen Whole Mouse Brain 10X: per-subclass discriminating "
            "markers plus GPCR specificity tiers.", False),
           ("168 genes covering 20 Allen subclass populations across the two "
            "regions; 77 of them are the minimum that still identifies all 20.", False),
           ("Code and audit trail on GitHub.", False)], accent=BLUE)

    callout(s, M, Inches(4.75), W - 2 * M, Inches(1.15),
            "The screen has told us WHERE and WHEN. It cannot tell us WHAT those neurons are, "
            "and it cannot tell us whether the behaviour needs them. Those are the two experiments "
            "in this deck.")
    footer(s, "Sources: TRAP_sample_manifest.csv; Fig06_7region_phase_delta.csv; Genelist_analysis_WMB v3.")


def s03_result(prs):
    s = blank(prs)
    title(s, "The result that drives the next phase",
          "Active−Passive TRAP density separates at Post and Reinstatement — not while the mice are taking morphine",
          kicker="Motivation")

    rows = [["Region", "Cluster", "Active − Passive\n(craving phases)", "Take forward?"],
            ["RE", "4", "1.19", ""],
            ["ORBm", "1", "1.04", "YES — Experiment 1 + 2"],
            ["CA", "1", "0.97", ""],
            ["LM", "4", "0.96", ""],
            ["AId", "1", "0.85", ""],
            ["BMAp", "4", "0.74", "YES — Experiment 1 + 2"],
            ["CP", "4", "0.31", ""]]
    hl = RGBColor(0xD8, 0xEC, 0xD8)
    cc = {(2, j): hl for j in range(4)}
    cc.update({(6, j): hl for j in range(4)})
    table(s, rows, M, BODY_TOP, Inches(6.0), col_w=[1.3, 0.9, 1.7, 2.1],
          align=["l", "c", "c", "l"], size=11, row_h=Inches(0.31), hdr_h=Inches(0.5),
          cell_colors=cc)

    bullets(s, [
        ("What the number is", 0, NAVY, True),
        ("Active minus Passive tdTomato+ density, z-scored within phase, "
         "averaged over Post and Reinstatement.", 1),
        ("What the phase pattern means", 0, NAVY, True),
        ("The difference is not largest during morphine taking. It is largest at Post "
         "and again at Reinstatement, and it narrows during Withdrawal.", 1),
        ("So the ensemble that distinguishes a mouse that earned its morphine from a "
         "mouse that merely received it is engaged at Post and re-engaged at relapse.", 1),
        ("Why ORBm and BMAp specifically", 0, NAVY, True),
        ("They come from the two different clusters, so together they sample a cortical "
         "and an amygdalar node rather than two versions of the same signal.", 1),
        ("Both have independent opioid literature: OFC is required for incubated oxycodone "
         "and heroin craving; amygdala ensembles are required for incubated craving across "
         "several drugs.", 1),
        ("Limitation I want to be explicit about", 0, RED, True),
        ("This is a correlation from a terminal design with 4-6 mice per phase. "
         "It is a hypothesis generator, not a causal claim. That is what Experiment 2 is for.", 1),
    ], COL2_L, BODY_TOP - Inches(0.05), COL2_W, BODY_H, size=12.5, gap=5)

    cap = s.shapes.add_textbox(M, Inches(4.55), Inches(6.0), Inches(0.3))
    cap.text_frame.text = "Where the Active−Passive difference appears across the 18-day paradigm"
    _tf(cap, 11, NAVY, bold=True)
    phase_timeline(s, M, Inches(4.9), Inches(6.0), [
        ("FR\n1-2", 2, NEUTRAL), ("Pre\n3-5", 3, WATER), ("During\n6-10", 5, MORPH),
        ("Post\n11-13", 3, MORPH), ("Withdr.\n14-16", 3, WATER), ("Re-exp\n17-18", 2, MORPH),
    ], height=Inches(0.5))

    gaps = [("", 0.0), ("gap opens", 11.5 / 18), ("gap narrows", 14.5 / 18), ("gap returns", 17.0 / 18)]
    for text, frac in gaps[1:]:
        tb = s.shapes.add_textbox(M + Inches(6.0 * frac) - Inches(0.6), Inches(5.48),
                                  Inches(1.2), Inches(0.42))
        tb.text_frame.word_wrap = True
        tb.text_frame.text = "\u25b2 " + text
        for para in tb.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.line_spacing = 0.85
            for r in para.runs:
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.color.rgb = GREEN if "narrow" not in text else AMBER
    legend(s, M, Inches(6.05), [("morphine", MORPH), ("water", WATER), ("fixed ratio", NEUTRAL)])
    footer(s, "Fig06_7region_phase_delta.csv; paired_slope_consistency.csv. Opioid literature cited on final slide.")


def s04_two_questions(prs):
    s = blank(prs)
    title(s, "Two questions, two experiments",
          "Experiment 1 gives us the target. Experiment 2 gives us the causal claim.",
          kicker="Plan")

    panel(s, M, BODY_TOP, COL2_W, Inches(3.55),
          "Q1  WHAT ARE THESE NEURONS?     →  EXPERIMENT 1: Xenium",
          [("Cell-type identity of the Post-tagged ensemble in ORBm and BMAp, "
            "in situ, at single-cell resolution.", False),
           ("Which of the 20 Allen subclasses the tdTomato+ cells actually belong to — "
            "is the Active-specific signal one population or many?", False),
           ("Which druggable GPCRs those specific cells express.", False),
           ("Whether the Post ensemble and the relapse ensemble are the same cells "
            "(tdTomato and Fos co-detected in the same section).", False),
           ("Output: a defensible molecular target, and the probe panel is already built.", True)],
          accent=BLUE, size=12.5)

    panel(s, COL2_L, BODY_TOP, COL2_W, Inches(3.55),
          "Q2  IS THE BEHAVIOUR NECESSARY?  →  EXPERIMENT 2: TRAP-chemogenetics",
          [("Silence the Post-tagged ensemble and ask whether morphine seeking falls.", False),
           ("Hypothesis: inhibiting the Post ensemble reduces seeking and relapse "
            "in Active mice, and does not in yoked Passive mice.", False),
           ("The Active-only prediction is the strong part of the hypothesis — "
            "it makes the experiment falsifiable rather than merely descriptive.", False),
           ("Output: a causal statement about a specific ensemble in a specific region.", True),
           ("This experiment needs one design change from my first sketch. Slides 8-9.", True)],
          accent=RGBColor(0x7A, 0x3E, 0x8C), size=12.5)

    callout(s, M, Inches(5.25), W - 2 * M, Inches(0.95),
            "Together they answer one sentence: the Post-active ensemble in region X consists of cell "
            "type Y expressing receptor Z, and relapse requires it. Either half alone is a weaker paper.",
            fill=LIGHT, edge=BLUE, size=13)


def s05_exp1_design(prs):
    s = blank(prs)
    title(s, "Experiment 1 — Xenium on the Post-tagged ensemble",
          "One design choice does most of the work: tag at Post, read out at relapse",
          kicker="Experiment 1")

    rows = [["", "Specification", "Why this and not something else"],
            ["Mice", "TRAP2;Ai14 (Fos2A-iCreER × Ai14)\nJAX 030323 × 007914",
             "Ai14 is the reporter you have chosen; tdTomato mRNA is the TRAP tag Xenium reads."],
            ["Groups", "Active n=4, Passive n=4",
             "At n=3 a rank-sum test cannot reach p<0.05 at all (best possible p = 0.10). "
             "n=4 makes p=0.029 attainable. Two mice buy the option of a real test."],
            ["Housing", "Single-housed from day 10\n(48 h before tagging), all groups alike",
             "BMAp abuts MEA, the social/conspecific nucleus. A cage mate inside the 6-h "
             "tagging window would tag social Fos and we would call it morphine."],
            ["Tag", "4-OHT 50 mg/kg i.p. immediately after the\nbehaviour day 12 session (Post day 2)",
             "TRAP2 labels within ~6 h centred on the injection, so this captures the Post state."],
            ["Behaviour", "Unchanged, days 1-18\n(FR, Pre, During, Post, Withdrawal, Re-exposure)",
             "Keeping the paradigm identical means the whole-brain TRAP result still maps onto it."],
            ["Abstinence", "Days 19-31, home cage",
             "Lets tdTomato accumulate and lets craving incubate."],
            ["Readout day", "Day 32: cue-induced seeking test,\nsacrifice 90 min later",
             "20 days after tagging: well past your 10-day minimum, and Fos from the test is at peak."],
            ["Tissue", "ORBm and BMAp coronal sections;\nActive and Passive on the SAME slide",
             "Puts the group comparison inside a slide instead of across slides."],
            ["Readouts", "tdTomato (Post ensemble) + Fos (relapse\nensemble) + 20 subclasses + GPCRs",
             "The overlap of tdTomato and Fos is the single most informative measurement here."]]
    table(s, rows, M, BODY_TOP, W - 2 * M, col_w=[0.95, 3.1, 4.4],
          align=["l", "l", "l"], size=10, hdr_size=11,
          row_h=Inches(0.45), hdr_h=Inches(0.28))

    callout(s, M, Inches(6.32), W - 2 * M, Inches(0.6),
            "Key point: tagging at Post and reading out at relapse turns one experiment into two "
            "measurements — what the Post ensemble is, and how much of it comes back at relapse.",
            fill=LIGHT, edge=BLUE, size=12)
    footer(s, "TRAP2 6-h window, 50 mg/kg 4-OHT and 48-h single housing before tagging: "
              "DeNardo et al. 2019 Nat Neurosci.")


def s06_exp1_panel(prs):
    s = blank(prs)
    title(s, "Experiment 1 — the panel fits, at half the add-on budget",
          "I matched our list against the real 10x base panel: the 77 priority-1 genes need only 51 of the 100 custom slots",
          kicker="Experiment 1")

    rows = [["Block", "What it answers", "Genes", "Free", "Add-on"],
            ["1  Cell-type unique", "Which of the 20 populations is this?", "62", "20", "42"],
            ["2  Cell-type GPCRs", "Druggable AND cell-type informative", "7", "1", "6"],
            ["4  IEG + reporter + class", "Active? TRAPped? Which class?", "8", "5", "3"],
            ["PRIORITY-1 TOTAL", "What we order: BMAp 38, ORBm 52", "77", "26", "51"],
            ["10x add-on cap", "Limit on a pre-designed panel", "—", "—", "100"],
            ["Slots still free", "Room for priority-2 backups", "—", "—", "49"]]
    cc = {(4, j): RGBColor(0xFF, 0xF2, 0xCC) for j in range(5)}
    cc.update({(6, j): RGBColor(0xD8, 0xEC, 0xD8) for j in range(5)})
    table(s, rows, M, BODY_TOP, Inches(7.15), col_w=[2.15, 2.9, 0.7, 0.65, 0.75],
          align=["l", "l", "c", "c", "c"], size=10.5, row_h=Inches(0.34),
          hdr_h=Inches(0.3), cell_colors=cc)

    panel(s, M, Inches(4.15), Inches(7.15), Inches(2.0),
          "THE ONLY PANEL RISK LEFT — TWO CUSTOM-SEQUENCE PROBES",
          [("tdTomato and iCre are not mouse genes, so they need advanced custom design "
            "from a FASTA of at least 80 bp, and 10x does not validate them.", True),
           ("tdTomato in Ai14 is very highly expressed. I will ask 10x to reduce its probe-set "
            "count, otherwise it crowds the optical field and depresses counts for every other gene.", False),
           ("If tdTomato still reads weakly: iCre reports the same cells independently, and "
            "tdTomato antibody staining on adjacent sections is a third, orthogonal route.", False)],
          accent=RGBColor(0x7A, 0x3E, 0x8C), size=11)

    panel(s, Inches(8.00), BODY_TOP, Inches(4.78), Inches(4.85),
          "WHAT THE CROSS-CHECK CHANGED",
          [("This was the open procurement risk in my first draft. It is now closed.", True),
           ("I pulled the 248-gene Xenium Mouse Brain v1 panel from 10x and matched it "
            "gene by gene against our list. 26 of the 77 priority-1 genes are already on it — "
            "including Fos and Arc, the two activity readouts — so we do not pay for them.", False),
           ("That leaves 51 add-on slots, half of the 100 available.", False),
           ("The 49 spare slots absorb 49 of the 85 priority-2 backups, so the panel we order "
            "can carry 146 of our 162 curated genes rather than 77.", False),
           ("Consequence for the budget: Xenium Prime 5K is no longer needed. "
            "Base Mouse Brain panel plus one add-on does everything, at the lower price point.", True),
           ("I would not fill all 100 blindly. I propose spending the spare slots on the "
            "priority-2 genes that rescue the subclasses which currently have no single unique "
            "marker, then on the GPCRs that already have ligands, and holding ~10 slots back.", False)],
          accent=GREEN, size=11)

    footer(s, "Base panel content read from the 10x gene_panel.json for Xenium_V1_FF_Mouse_Brain_MultiSection_1 "
              "(248 targets + 27 negative controls). Cap and exogenous-target rules: 10x CG000643 / CG000683. "
              "Full match table: outputs/Xenium_addon_vs_base_panel.csv")


def s07_exp1_schedule(prs):
    s = blank(prs)
    title(s, "Experiment 1 — schedule",
          "The September pilot deliberately uses the base panel only, so tissue QC is not blocked by panel design",
          kicker="Experiment 1")

    rows = [["When", "Step", "What must be true to proceed"],
            ["Sep 2026", "Xenium pilot with Alina's lab. Base mouse brain panel only, "
                         "1 Active + 1 Passive brain from existing tissue.",
             "Nothing — this is the gate, and it needs no new panel."],
            ["Sep 2026", "QC criteria fixed in advance: transcripts per cell, "
                         "negative-probe rate, segmentation quality, registration of "
                         "ORBm and BMAp to Allen, and recovery of the expected subclasses.",
             "Agree the numeric thresholds with Alina BEFORE the run, so pass/fail is not a judgement call."],
            ["Sep-Oct 2026", "Choose which priority-2 genes fill the 49 spare add-on slots. "
                             "The cross-check against the base panel is already done.",
             "Pilot QC passed. This is now a scientific choice, not a feasibility question."],
            ["Oct 2026", "Submit the design: base Mouse Brain panel + ~85-gene add-on, "
                         "including FASTA for tdTomato and iCre.",
             "Ai14 tdTomato and Fos2A-iCreER sequences in hand — request these now, they are "
             "the one item with no substitute."],
            ["Oct-Nov 2026", "Panel manufacture and delivery (allow 6-8 weeks).",
             "Order placed."],
            ["Nov 2026", "Run the behaviour cohort for Experiment 1 (n=4+4) "
                         "so tissue is ready when the panel arrives.",
             "Mice genotyped TRAP2;Ai14; single-housed from day 10; 4-OHT prepared fresh."],
            ["Dec 2026 - Jan 2027", "Xenium run on ORBm and BMAp; analysis.",
             "Panel delivered; sections cut and QC'd."]]
    table(s, rows, M, BODY_TOP, W - 2 * M, col_w=[1.35, 4.6, 3.9],
          align=["l", "l", "l"], size=10.5, row_h=Inches(0.62), hdr_h=Inches(0.3))

    callout(s, M, Inches(6.35), W - 2 * M, Inches(0.6),
            "The critical path is now the two transgene sequences, not the gene list. "
            "Advanced custom design has a longer queue than a standard add-on, so the FASTA files "
            "need to be in hand before we submit in October.",
            size=12)


def s08_problem(prs):
    s = blank(prs)
    title(s, "Experiment 2 — the problem with my first sketch",
          "Injecting CNO on Post day 7 would have tested nothing, because the receptor is not there yet",
          kicker="Experiment 2  ·  the problem")

    panel(s, M, BODY_TOP, COL2_W, Inches(2.35),
          "TRAP IS INTRINSICALLY RETROSPECTIVE",
          [("To be tagged, a neuron must be active. To be silenced, it must first "
            "manufacture the receptor. Those cannot happen at the same time.", True),
           ("So TRAP can never answer 'does activity during Post cause behaviour "
            "during Post'. It answers 'are the neurons that were active during Post "
            "required for behaviour later'.", False),
           ("This is a property of the method, not a flaw in the plan — "
            "but it decides what the experiment is allowed to claim.", False)],
          accent=RED, size=12.5)

    panel(s, COL2_L, BODY_TOP, COL2_W, Inches(2.35),
          "AND THE ARITHMETIC DOES NOT WORK",
          [("AAV-DIO-hM4Di needs about 14 days after 4-OHT before hM4Di is "
            "functional, on top of ~3 weeks for the AAV itself before tagging.", True),
           ("My sketch: 4-OHT on Post day 2 (behaviour day 12), CNO on 'Post day 7' "
            "— five days later.", False),
           ("Five days is not enough. We would inhibit nothing, see no behavioural "
            "effect, and be unable to tell that apart from the hypothesis being wrong.", False)],
          accent=RED, size=12.5)

    rows = [["Requirement", "Needs", "My sketch gave it", "Verdict"],
            ["AAV established before tagging", "~21 days", "not specified", "fixable"],
            ["hM4Di functional after 4-OHT", "~14 days", "5 days", "FAILS"],
            ["Post phase long enough to test inside it", "~17 days", "3 days (day 11-13)", "FAILS"]]
    cc = {(2, 3): RGBColor(0xFD, 0xD8, 0xD8), (3, 3): RGBColor(0xFD, 0xD8, 0xD8),
          (1, 3): RGBColor(0xFF, 0xF2, 0xCC)}
    table(s, rows, M, Inches(3.92), Inches(8.4), col_w=[3.4, 1.2, 1.6, 1.1],
          align=["l", "c", "c", "c"], size=11, row_h=Inches(0.33), hdr_h=Inches(0.32),
          cell_colors=cc)

    callout(s, Inches(9.25), Inches(3.92), Inches(3.53), Inches(1.65),
            "A null result from an underpowered, under-expressed experiment is the "
            "worst outcome available to us: it costs a full cohort and answers nothing. "
            "The next slide is the fix.",
            fill=RGBColor(0xFD, 0xE9, 0xE9), edge=RED, size=12)
    footer(s, "14-day post-4-OHT window for AAV-delivered DREADDs: TRAP2 ghrelin study, Int J Mol Sci 2022. "
              "3-4 weeks for Cre-dependent AAV expression is standard.")


def s09_solution(prs):
    s = blank(prs)
    title(s, "Experiment 2 — the constraint points at the right experiment",
          "Three independent timescales all land on the same window, about two weeks after the last morphine",
          kicker="Experiment 2  ·  the fix")

    rows = [["Timescale", "Value", "Source"],
            ["hM4Di becomes functional after 4-OHT", "~14 days", "TRAP2 + AAV-DIO-DREADD studies"],
            ["Incubation of craving peaks", "~withdrawal day 14", "Fos ensemble studies, inverted-U across drugs"],
            ["OFC inactivation reduces opioid seeking", "at day 15, NOT at day 1", "Oxycodone incubation, OFC inactivation"]]
    table(s, rows, M, BODY_TOP, Inches(7.4), col_w=[3.3, 1.9, 2.9],
          align=["l", "c", "l"], size=11, row_h=Inches(0.42), hdr_h=Inches(0.3))

    callout(s, M, Inches(3.05), Inches(7.4), Inches(1.15),
            "The technique is ready exactly when the behaviour is maximal, and exactly when "
            "the region is known to matter. The delay we were fighting is the delay we want.",
            fill=RGBColor(0xD8, 0xEC, 0xD8), edge=GREEN, size=13)

    panel(s, M, Inches(4.42), Inches(7.4), Inches(1.75),
          "WHAT CHANGES, AND WHAT DOES NOT",
          [("Does NOT change: behaviour days 1-18 stay exactly as run. Same FR, Pre, During, "
            "Post, Withdrawal, Re-exposure. The whole-brain TRAP and Xenium data still map on.", True),
           ("Changes: we add a home-cage abstinence period and a drug test at the end. "
            "We add days after day 18; we do not alter days 1-18.", True),
           ("That is the whole argument for this redesign — it is additive, not a new paradigm.", False)],
          accent=GREEN, size=12)

    panel(s, Inches(8.25), BODY_TOP, Inches(4.53), Inches(3.85),
          "THE QUESTION, RESTATED HONESTLY",
          [("Cannot test: does Post activity cause Post seeking?", True),
           ("Can test, and this was the second half of my own hypothesis:", False),
           ("Are ORBm / BMAp neurons that were active during Post necessary for "
            "morphine seeking after protracted abstinence, and is that requirement "
            "specific to mice that earned the drug?", True),
           ("This is also the more clinically relevant question. Patients do not "
            "relapse during use; they relapse after abstinence.", False)],
          accent=GREEN, size=12.5)

    panel(s, Inches(8.25), Inches(4.78), Inches(4.53), Inches(0.82),
          "IF YOU WANT THE WITHIN-POST TEST ANYWAY",
          [("There is a way to do it at 48 h instead of 14 days. Slide 13.", False)],
          accent=AMBER, size=12)
    footer(s, "Full citations on the final slide.")


def s10_exp2_design(prs):
    s = blank(prs)
    title(s, "Experiment 2 — design",
          "2 × 2 between subjects, drug crossed within subject; one region at a time",
          kicker="Experiment 2")

    rows = [["", "Levels", "Purpose"],
            ["Group (between)", "Active  /  yoked Passive",
             "The Active-only prediction. Passive is the specificity control."],
            ["Virus (between)", "AAV-DIO-hM4Di-mCherry  /  AAV-DIO-mCherry",
             "The mCherry + drug arm is the essential control: it removes any effect of the "
             "ligand itself. More important than a saline arm."],
            ["Drug (WITHIN subject)", "DCZ 0.1 mg/kg  /  vehicle, crossover",
             "Each mouse is its own control, which roughly halves the mice needed."],
            ["Region", "ORBm first, then BMAp",
             "ORBm has the larger Active−Passive difference and is the easier bilateral target. "
             "BMAp second, informed by the Xenium result."]]
    table(s, rows, M, BODY_TOP, Inches(7.5), col_w=[1.6, 2.4, 4.1],
          align=["l", "l", "l"], size=10.5, row_h=Inches(0.62), hdr_h=Inches(0.28))

    prows = [["n per cell", "dz=0.8", "dz=1.0", "dz=1.2", "dz=1.5"],
             ["4  (my sketch)", "21%", "29%", "38%", "53%"],
             ["8", "50%", "68%", "83%", "95%"],
             ["10", "62%", "80%", "92%", "99%"],
             ["12", "71%", "88%", "97%", "100%"]]
    cc = {(1, j): RGBColor(0xFD, 0xD8, 0xD8) for j in range(5)}
    cc.update({(3, j): RGBColor(0xD8, 0xEC, 0xD8) for j in range(5)})
    table(s, prows, Inches(8.35), BODY_TOP, Inches(4.43), col_w=[1.5, 0.9, 0.9, 0.9, 0.9],
          align=["l", "c", "c", "c", "c"], size=10.5, row_h=Inches(0.3),
          hdr_h=Inches(0.3), cell_colors=cc)

    bullets(s, [
        ("Power, paired t-test, two-tailed α=0.05", 0, NAVY, True),
        ("Reported chemogenetic effects on drug seeking are large, roughly a "
         "40-50% reduction, so dz ~ 1.0-1.2 is a fair planning assumption.", 1),
        ("n=4 detects nothing smaller than dz~1.6. n=10 is the first row where a "
         "large effect is reliably detectable.", 1),
        ("What I propose instead of choosing one n", 0, NAVY, True),
        ("Cohort 1 = expression pilot, n=4 per cell (16 mice, ORBm). "
         "No behavioural inference claimed. It buys the AAV coordinates, the "
         "percentage of tdTomato+ cells that become hM4Di+, confirmation that DCZ "
         "suppresses Fos in those cells, and a variance estimate.", 1),
        ("Cohort 2 = powered, n=10 per cell (40 mice, ORBm), with n set by "
         "cohort 1's variance rather than by my guess.", 1),
        ("Allow ~25% attrition for missed targeting and health, so order accordingly.", 1),
    ], Inches(8.35), Inches(3.05), Inches(4.43), Inches(3.4), size=11.5, gap=5)

    callout(s, M, Inches(4.53), Inches(7.5), Inches(1.05),
            "Starting coordinates to verify in cohort 1, then fix: ORBm ~ AP +2.6, ML ±0.35, DV −2.7; "
            "BMAp ~ AP −1.8, ML ±2.4, DV −5.1. Use 150-250 nl at low rate — BMAp sits next to MEA, BLA and "
            "the sAMY-like GABA populations our marker work flagged, so spread must be verified histologically.",
            size=11, bold=False)


def s11_exp2_schedule(prs):
    s = blank(prs)
    title(s, "Experiment 2 — day-by-day schedule",
          "Behaviour days 1-18 are unchanged; everything new is appended after day 18",
          kicker="Experiment 2")

    rows = [["Day", "Phase / step", "Reward", "Note"],
            ["−14", "Bilateral AAV surgery (hM4Di or mCherry), ORBm", "—",
             "26 days before tagging, comfortably past the ~21-day AAV requirement"],
            ["−13 to 0", "Recovery, handling, injection habituation", "—",
             "Handling matters: it lowers TRAP background labelling"],
            ["1-2", "Fixed ratio", "—", "As currently run"],
            ["3-5", "Pre", "Water", "As currently run"],
            ["6-10", "During", "Morphine", "Passive lick lockout applies, as currently run"],
            ["11-13", "Post", "Morphine", "4-OHT 50 mg/kg i.p. right after the DAY 12 session"],
            ["14-16", "Withdrawal", "Water", "As currently run"],
            ["17-18", "Re-exposure", "Morphine", "As currently run. Day 18 = last morphine"],
            ["19-31", "Forced abstinence, home cage", "—",
             "hM4Di matures (19 days post-4-OHT by day 31) and craving incubates"],
            ["32", "TEST 1 — cue-induced seeking, no morphine", "—",
             "DCZ or vehicle 15 min before. Withdrawal day 14. Primary endpoint"],
            ["34", "TEST 2 — same test, treatment crossed over", "—",
             "48 h washout is ample for DCZ"],
            ["36", "SPECIFICITY — water or sucrose PR under DCZ", "Water",
             "Shows the effect is not general motivation"],
            ["38", "MOTOR — locomotion / open field under DCZ", "—",
             "Rules out a movement confound"],
            ["40", "Final session, sacrifice 90 min later", "—",
             "Histology: placement, tdTomato ∩ hM4Di overlap, Fos suppression"]]
    cc = {}
    for i, r in enumerate(rows):
        if str(r[1]).startswith(("TEST", "SPECIFICITY", "MOTOR")):
            cc.update({(i, j): RGBColor(0xD8, 0xEC, 0xD8) for j in range(4)})
    cc.update({(6, j): RGBColor(0xFF, 0xF2, 0xCC) for j in range(4)})
    table(s, rows, M, BODY_TOP, Inches(7.55), col_w=[0.7, 2.6, 0.75, 3.5],
          align=["c", "l", "c", "l"], size=9, hdr_size=9.5,
          row_h=Inches(0.30), hdr_h=Inches(0.24), cell_colors=cc)

    cap = s.shapes.add_textbox(Inches(8.40), BODY_TOP - Inches(0.02), Inches(4.38), Inches(0.28))
    cap.text_frame.text = "Same picture as one timeline"
    _tf(cap, 12, NAVY, bold=True)
    # Units are days except the AAV block, which compresses days -14..0 into 4
    # units so the behavioural phases stay legible. Marker fractions are derived
    # from the same cumulative units rather than eyeballed.
    tl_segs = [
        ("AAV\n−14", 4, NAVY), ("FR\n1-2", 2, NEUTRAL), ("Pre", 3, WATER),
        ("During", 5, MORPH), ("Post", 3, MORPH), ("W/D", 3, WATER), ("Re", 2, MORPH),
        ("abstinence 19-31", 13, ABST), ("tests\n32-40", 9, GREEN),
    ]
    TL_L, TL_W = Inches(8.40), Inches(4.38)
    phase_timeline(s, TL_L, Inches(2.35), TL_W, tl_segs, height=Inches(0.62), label_size=8)

    units = sum(sg[1] for sg in tl_segs)
    cum = {}
    acc = 0
    for label, d, _ in tl_segs:
        cum[label] = (acc, acc + d)
        acc += d
    post_a, post_b = cum["Post"]
    arrow_specs = [
        ((post_a + (post_b - post_a) / 2) / units, "4-OHT\nday 12", RED),
        (cum["tests\n32-40"][0] / units, "DCZ test 1\nday 32", GREEN),
    ]
    for frac, text, col in arrow_specs:
        cx = TL_L + Emu(int(TL_W * frac))
        # Draw the pin below the bar only; running it through the bar cut the
        # phase label in half.
        pin = s.shapes.add_shape(1, cx - Emu(11000), Inches(2.97), Emu(22000), Inches(0.13))
        pin.fill.solid()
        pin.fill.fore_color.rgb = col
        pin.line.fill.background()
        pin.shadow.inherit = False
        lab = s.shapes.add_textbox(cx - Inches(0.7), Inches(3.1), Inches(1.4), Inches(0.42))
        lab.text_frame.word_wrap = True
        lab.text_frame.text = text
        for para in lab.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.line_spacing = 0.85
            for r in para.runs:
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.color.rgb = col

    panel(s, Inches(8.40), Inches(3.65), Inches(4.38), Inches(2.4),
          "THE 20-DAY GAP IS THE POINT",
          [("Day 12 (4-OHT) to day 32 (first test) = 20 days.", True),
           ("hM4Di needs ~14 days after tagging  →  satisfied with 6 days of margin.", False),
           ("Your own requirement of at least 10 days after TRAP  →  satisfied.", False),
           ("Day 32 is withdrawal day 14 from the last morphine, which is where "
            "incubated opioid seeking peaks and where OFC inactivation is known to work.", False),
           ("Days 1-18 are untouched, so the whole-brain TRAP and Xenium results "
            "still describe this cohort.", True)],
          accent=GREEN, size=11)

    legend(s, Inches(8.40), Inches(6.15),
           [("morphine", MORPH), ("water", WATER), ("abstinence", ABST), ("test", GREEN)])
    footer(s, "Primary endpoints: active licks and PR breakpoint in the 1-h seeking test. "
              "Analysis: Group × Virus × Drug mixed-effects model with mouse as random intercept.")


def s12_controls(prs):
    s = blank(prs)
    title(s, "Experiment 2 — the controls a reviewer will ask for",
          "Each one closes a specific alternative explanation",
          kicker="Experiment 2")

    rows = [["Control", "What it rules out", "Cost"],
            ["AAV-DIO-mCherry + DCZ", "That the ligand itself changed behaviour. This is the "
             "control that matters most and it is why I am not relying on a saline arm.",
             "Built into the 2×2"],
            ["Vehicle in hM4Di mice (crossover)", "Baseline drift and order effects.", "Free — within subject"],
            ["No-4-OHT (oil vehicle) + AAV", "Cre-independent recombination leak, i.e. hM4Di in "
             "cells that were never active.", "3-4 mice, histology only"],
            ["Home-cage 4-OHT, no session", "That we tagged the paradigm rather than the Post state. "
             "Gives the background labelling rate to subtract.", "3-4 mice, histology only"],
            ["Water or sucrose PR under DCZ", "A general drop in motivation or in licking ability.", "One session, same mice"],
            ["Locomotion under DCZ", "A motor confound.", "One session, same mice"],
            ["Fos in hM4Di+ cells after DCZ", "That we inhibited nothing. This is target engagement "
             "and without it a null result is uninterpretable.", "Terminal histology"]]
    table(s, rows, M, BODY_TOP, Inches(8.4), col_w=[2.5, 4.6, 1.6],
          align=["l", "l", "l"], size=10.5, row_h=Inches(0.51), hdr_h=Inches(0.28))

    panel(s, Inches(9.25), BODY_TOP, Inches(3.53), Inches(2.6),
          "USE DCZ, NOT CNO",
          [("CNO back-converts to clozapine in vivo, and clozapine crosses the "
            "blood-brain barrier and alters motivation on its own.", True),
           ("In a drug-seeking assay that confound sits directly on top of the "
            "dependent variable, so it is not a risk we should take.", False),
           ("DCZ: higher affinity, faster brain entry, no detectable activity across "
            "318 off-target GPCRs. Effective at 0.1 mg/kg i.p. or lower.", False)],
          accent=RGBColor(0x7A, 0x3E, 0x8C), size=11.5)

    panel(s, Inches(9.25), Inches(4.22), Inches(3.53), Inches(2.5),
          "HOUSING — RESOLVED, SINGLE CAGE",
          [("I had flagged a possible conflict with the yoked design. There is none: "
            "yoking is implemented in software, by replaying the Active partner's reward "
            "log into a separate Passive session, so partners never share a box.", False),
           ("So we single-house from day 10 at no cost to the paradigm — and the paradigm "
            "runs in the operant rig, not the home cage, so the apparatus is unchanged too.", False),
           ("Applied identically to Active, Passive and control, and started 48 h ahead so "
            "the isolation itself is not the salient event being tagged.", True)],
          accent=GREEN, size=10.5)
    footer(s, "CNO reverse metabolism: Gomez et al. 2017 Science; Manvich et al. 2018 Sci Rep. "
              "DCZ: Nagai et al. 2020 Nat Neurosci. Yoking implementation: MicroLogReplayExperiment, "
              "experiments/ratio_tasks.py.")


def s13_alternatives(prs):
    s = blank(prs)
    title(s, "If you want inhibition inside the Post window",
          "There is a clean way to do it, but it needs new breeding — so I propose it as phase 2",
          kicker="Experiment 2  ·  alternatives")

    rows = [["", "A — Protracted abstinence\n(RECOMMENDED, start Nov)",
             "B — Transgenic hM4Di + cannula\n(phase 2)",
             "C — Extended Post\n(not recommended)"],
            ["Effector", "AAV-DIO-hM4Di into ORBm / BMAp",
             "TRAP2 × R26-LSL-hM4Di (JAX 026219).\nEffector already in the genome, so no AAV lag",
             "AAV-DIO-hM4Di"],
            ["Delay from 4-OHT to test", "~20 days", "48 h",
             "~15 days, but inside an extended Post"],
            ["Region specificity", "From the injection site",
             "From local DCZ infusion through cannulae —\nthe transgene is brain-wide",
             "From the injection site"],
            ["What it tests", "Is the Post ensemble needed for relapse\nafter abstinence?",
             "Is the Post ensemble needed for seeking\nwithin the Post phase?",
             "Same as B, but in a paradigm we\nhave not characterised"],
            ["Main cost", "None beyond the appended days",
             "New cross; Rosa26 conflict means this line\ncannot also carry Ai14, so the tag becomes\nmCitrine. Cannula surgery. 4-6 months of breeding",
             "3 extra weeks of morphine per mouse,\nso drug history no longer matches the\n18-day model behind our TRAP data"],
            ["Ready to start", "November 2026", "Mid 2027 if we start breeding now", "November 2026"]]
    cc = {}
    for i in range(len(rows)):
        cc[(i, 1)] = RGBColor(0xD8, 0xEC, 0xD8) if i else GREEN
        if i:
            cc[(i, 3)] = RGBColor(0xF2, 0xF2, 0xF2)
    table(s, rows, M, BODY_TOP, W - 2 * M, col_w=[1.45, 3.0, 4.1, 3.0],
          align=["l", "l", "l", "l"], size=9.5, hdr_size=10.5,
          row_h=Inches(0.68), hdr_h=Inches(0.5), cell_colors=cc)

    callout(s, M, Inches(6.35), W - 2 * M, Inches(0.62),
            "My recommendation: run A now, and start the B cross in parallel so the within-Post question "
            "is answerable next year. A and B together are a much stronger story than either alone, "
            "because they separate 'needed while taking' from 'needed for relapse'.",
            fill=RGBColor(0xD8, 0xEC, 0xD8), edge=GREEN, size=12)
    footer(s, "48-h transgenic hM4Di precedent with intracerebral CNO: LEC engram study, bioRxiv 2025. "
              "R26-LSL-hM4Di: JAX 026219, Rosa26 locus — same locus as Ai14 (JAX 007914).")


def s14_gantt(prs):
    s = blank(prs)
    title(s, "Master timeline", "September 2026 to August 2027", kicker="Schedule")

    months = ["Sep 26", "Oct", "Nov", "Dec", "Jan 27", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug"]
    left0, top0 = M + Inches(3.1), BODY_TOP + Inches(0.05)
    colw = (W - 2 * M - Inches(3.1)) / len(months)
    ROW = Inches(0.335)

    for i, m in enumerate(months):
        b = s.shapes.add_textbox(left0 + colw * i, BODY_TOP - Inches(0.28), colw, Inches(0.3))
        b.text_frame.text = m
        _tf(b, 9.5, NAVY, bold=True, line=0.9)
        b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    tasks = [
        ("Xenium pilot with Alina (base panel)", 0, 1, BLUE),
        ("Panel cross-check + quote", 0, 2, RGBColor(0x8A, 0xB4, 0xD8)),
        ("Order custom panel", 1, 1, BLUE),
        ("Panel manufacture", 1, 3, RGBColor(0x8A, 0xB4, 0xD8)),
        ("Exp 1 behaviour cohort (n=3+3)", 2, 2, BLUE),
        ("Exp 1 Xenium run + analysis", 3, 3, BLUE),
        ("Exp 2 cohort 1 surgery (ORBm, n=16)", 2, 1, RGBColor(0x7A, 0x3E, 0x8C)),
        ("Exp 2 cohort 1 behaviour + DCZ tests", 2, 3, RGBColor(0x9D, 0x6A, 0xB0)),
        ("Cohort 1 histology + power estimate", 4, 2, RGBColor(0x9D, 0x6A, 0xB0)),
        ("DECISION GATE: go / no-go on powered cohort", 5, 1, AMBER),
        ("Exp 2 cohort 2 powered (ORBm, n=40)", 5, 4, RGBColor(0x7A, 0x3E, 0x8C)),
        ("Exp 2 BMAp cohort", 8, 4, RGBColor(0x7A, 0x3E, 0x8C)),
        ("Start TRAP2 × R26-hM4Di cross (option B)", 2, 8, GREEN),
        ("Writing / figures", 9, 3, GREY),
    ]
    for i in range(len(months) + 1):
        ln = s.shapes.add_shape(1, left0 + colw * i, top0 - Inches(0.04), Emu(6350),
                                ROW * len(tasks) + Inches(0.04))
        ln.fill.solid()
        ln.fill.fore_color.rgb = BAND
        ln.line.fill.background()
        ln.shadow.inherit = False

    for i, (name, start, dur, col) in enumerate(tasks):
        y = top0 + ROW * i
        lbl = s.shapes.add_textbox(M, y - Inches(0.03), Inches(3.05), Inches(0.28))
        lbl.text_frame.text = name
        _tf(lbl, 9.5, NAVY if col != GREY else GREY, bold=col == AMBER, line=0.9)
        bar = s.shapes.add_shape(5, left0 + colw * start + Emu(9000), y,
                                 colw * dur - Emu(18000), Inches(0.225))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        bar.shadow.inherit = False

    callout(s, M, top0 + ROW * len(tasks) + Inches(0.14), W - 2 * M, Inches(0.55),
            "The decision gate in February is deliberate: no powered cohort is ordered until cohort 1 has "
            "shown that hM4Di reaches the tagged cells and that DCZ suppresses their activity.",
            size=11.5)


def s15_decisions(prs):
    s = blank(prs)
    title(s, "What I need from you, and what could still go wrong",
          kicker="Decisions")

    panel(s, M, BODY_TOP, COL2_W, Inches(4.6), "DECISIONS I NEED",
          [("1.  Which region first for chemogenetics. I recommend ORBm — larger effect in "
            "the screen, easier bilateral target, and cortex is more forgiving of spread.", False),
           ("2.  Group sizes: do you accept cohort 1 as an explicitly underpowered "
            "expression pilot, with the powered cohort gated on its result?", False),
           ("3.  Sex. The behaviour cohort was 7 female / 7 male with within-sex yoked pairs, "
            "but sex is not recorded for the first 8 TRAP mice. I would balance sex in the "
            "powered cohort rather than in the pilot.", False),
           ("4.  Whether to start the TRAP2 × R26-hM4Di cross now, so the within-Post "
            "question is answerable in 2027.", False),
           ("5.  Two records I could not find anywhere in the code, manifests or notes, "
            "and need from the lab: the housing condition and the exact 4-OHT protocol "
            "used for the existing 20-mouse TRAP cohort.", True)],
          accent=BLUE, size=11.5)

    panel(s, COL2_L, BODY_TOP, COL2_W, Inches(4.6), "RISKS, AND WHAT I WOULD DO",
          [("If the existing TRAP cohort was group-housed, its labelling baseline differs from "
            "the single-housed cohort, and the Xenium result is not strictly comparable to the "
            "density screen that chose these two regions. Mitigation: the home-cage 4-OHT "
            "control measures the new baseline directly, so the comparison can be corrected "
            "rather than lost.", True),
           ("Custom transgene probes are not validated by 10x. If tdTomato detection is weak, "
            "iCre gives a second independent readout of the same cells, and immunostaining on "
            "adjacent sections is the fallback.", False),
           ("BMAp is small and next to MEA, BLA and the sAMY-like GABA populations. "
            "Mitigation: small injection volume, and every mouse verified histologically against "
            "the subclass markers from the panel work.", False),
           ("Our Post ensemble may not be one cell type. If Xenium shows it is heterogeneous, "
            "the chemogenetic result becomes harder to interpret — which is a reason to let "
            "Experiment 1 report before the powered cohort is locked.", False),
           ("Morphine dose is now specified (next slide). The one number still missing is the "
            "Micro4 flow rate: the rig pulses the pump for 0.3 s per reward, so volume per "
            "reward = rate × 0.3 s. One calibration measurement closes the methods section.", False)],
          accent=AMBER, size=11.5)


def s16_params(prs):
    s = blank(prs)
    title(s, "Parameters now locked",
          "Nothing on this slide is still open, so the October panel order and the November surgery can both proceed",
          kicker="Locked")

    rows = [["", "Value", "Note"],
            ["Driver × reporter", "TRAP2 (Fos2A-iCreER, JAX 030323) × Ai14 (JAX 007914)",
             "tdTomato mRNA is the permanent tag that Xenium reads."],
            ["Tagging", "4-OHT 50 mg/kg i.p., single dose, straight after the day-12 session",
             "Labels within a ~6 h window centred on the injection."],
            ["Housing", "Single cage from day 10 onward, applied identically to every group",
             "48 h of habituation before tagging. Costs nothing: yoking is software, "
             "and the task runs in the rig rather than the home cage."],
            ["Xenium panel", "Base Mouse Brain panel (248 genes) + custom add-on",
             "77 priority-1 genes need 51 add-on slots; 49 slots spare for priority-2 backups."],
            ["Custom probes", "tdTomato, iCre",
             "The only two non-mouse targets, and the only long-lead item."],
            ["Xenium group size", "4 Active + 4 Passive per region, plus 2 tagging controls",
             "Primary analyses are within-animal, so this n is doing more work than it looks."],
            ["Chemogenetic ligand", "DCZ, 0.1 mg/kg i.p.",
             "Chosen over CNO, which back-converts to clozapine — see slide 12."],
            ["Morphine, target", "15 mg/kg/day, self-titrated by licking",
             "Dose is set by how many rewards the animal earns, not by the experimenter."],
            ["Morphine, measured", "~16 mg/kg at 30 rewards (PR 30)  →  ~0.53 mg/kg per reward",
             "~13 µg per reward in a 25 g mouse. Volume per reward follows from the "
             "Micro4 flow rate × the 0.3 s pulse the rig sends."]]
    table(s, rows, M, BODY_TOP, W - 2 * M, col_w=[1.85, 4.35, 6.0],
          align=["l", "l", "l"], size=10, hdr_size=10.5,
          row_h=Inches(0.45), hdr_h=Inches(0.26))

    panel(s, M, Inches(6.05), W - 2 * M, Inches(1.25),
          "WHY A SELF-TITRATED DOSE IS AN ADVANTAGE HERE, NOT A WEAKNESS",
          [("Active and Passive intake are matched by construction rather than by assumption: "
            "the Passive session replays the Active partner's reward log, so reward count and "
            "timing are identical. That is precisely what makes the Active-versus-Passive "
            "contrast a test of volition rather than a test of drug exposure.", True),
           ("I will still report mg/kg per mouse per session and carry it as a covariate, "
            "because reviewers of oral self-administration always ask for it.", False)],
          accent=BLUE, size=11)


def s17_refs(prs):
    s = blank(prs)
    title(s, "References for every parameter used in this deck", kicker="References")

    rows = [["Claim in this deck", "Source"],
            ["TRAP2 labels within a ~6 h window centred on 4-OHT; 4-OHT 50 mg/kg i.p. at "
             "10 mg/mL in 1:4 castor:sunflower oil; single housing 48 h before tagging",
             "DeNardo LA et al. Temporal evolution of cortical ensembles promoting remote memory "
             "retrieval. Nat Neurosci 2019;22:460-469"],
            ["TRAP2 = Fos2A-iCreER, JAX 030323;  Ai14 = JAX 007914",
             "DeNardo et al. 2019; Madisen L et al. Nat Neurosci 2010;13:133-140"],
            ["AAV-delivered Cre-dependent DREADDs need ~2 weeks after 4-OHT, on top of "
             "~3 weeks for AAV expression before tagging",
             "López-Ferreras L et al. TRAPing ghrelin-activated circuits. Int J Mol Sci 2022;23:559; "
             "Cre-dependent AAV DREADD studies use 3-4 week expression windows"],
            ["Transgenic Cre-dependent hM4Di can be used 48 h after tagging, with local "
             "intracerebral ligand infusion for region specificity",
             "LEC engram cell study, bioRxiv 2025 (TRAP2 × hM4Di-DREADD, 48 h, intra-LEC CNO)"],
            ["R26-LSL-hM4Di/mCitrine = JAX 026219, targeted to the Rosa26 locus "
             "(hence the conflict with Ai14)",
             "Zhu H et al., Cre-dependent DREADD mice; JAX strain 026219; MGI:5614050"],
            ["CNO is reverse-metabolised to clozapine and produces clozapine-like effects",
             "Gomez JL et al. Science 2017;357:503-507; Manvich DF et al. Sci Rep 2018;8:3840"],
            ["DCZ is more potent and selective than CNO; effective at 0.001-0.1 mg/kg with no "
             "activity across 318 off-target GPCRs",
             "Nagai Y et al. Nat Neurosci 2020;23:1157-1167; Nentwig TB et al. Sci Rep 2022;12:6520"],
            ["OFC inactivation reduces oxycodone seeking on withdrawal day 15 but not day 1",
             "Fredriksson I / Shaham lab. Role of orbitofrontal cortex in incubation of oxycodone "
             "craving. Addict Biol 2020;25:e12927"],
            ["Fos ensembles are causally required for opioid and nicotine seeking; incubation "
             "peaks around withdrawal day 14",
             "Funk D et al. J Neurosci 2016 (CeA, nicotine); Fredriksson I et al. Sci Adv 2023 "
             "(vSub, oxycodone); Warren BL et al. Addict Biol 2022 (IL, oxycodone)"],
            ["Xenium supports exogenous targets (reporters, transgenes) via advanced custom "
             "design from a FASTA of ≥80 bp, but 10x does not validate them; add-on panels are "
             "capped at 100 genes",
             "10x Genomics CG000683 Xenium Custom Panel Design; CG000643 Designing Custom "
             "Xenium Panels"],
            ["Xenium Mouse Brain v1 base panel = 248 genes (+27 negative controls); 26 of our "
             "77 priority-1 genes are already on it, so the add-on requirement is 51",
             "10x gene_panel.json for Xenium_V1_FF_Mouse_Brain_MultiSection_1; matched in "
             "v3/E_Planning against the D07 ordering list. Table: outputs/Xenium_addon_vs_base_panel.csv"],
            ["At n=3 vs 3 the smallest attainable two-tailed rank-sum p is 0.10, so p<0.05 is "
             "unreachable; at n=4 vs 4 it is 0.029",
             "Exact Mann-Whitney null distribution: 20 and 70 possible orderings respectively. "
             "Verified in v3/E_Planning power checks"]]
    table(s, rows, M, BODY_TOP, W - 2 * M, col_w=[4.6, 7.6],
          align=["l", "l"], size=9, hdr_size=10.5, row_h=Inches(0.42), hdr_h=Inches(0.24))


def main() -> None:
    prs = deck()
    for fn in (s01_title, s02_where_we_are, s03_result, s04_two_questions,
               s05_exp1_design, s06_exp1_panel, s07_exp1_schedule, s08_problem,
               s09_solution, s10_exp2_design, s11_exp2_schedule, s12_controls,
               s13_alternatives, s14_gantt, s15_decisions, s16_params, s17_refs):
        fn(prs)
    out = Path(__file__).resolve().parents[2] / "outputs" / "NextPhase_Plan_ORBm_BMAp_TRAP.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"[DONE] {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
