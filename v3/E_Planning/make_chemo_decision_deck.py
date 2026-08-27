#!/usr/bin/env python
"""
PI-facing chemogenetics decision deck.

One figure per slide. Title + one-line message only.
Options are laid out so a PI can pick A/B/C/D/E without reading a protocol.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1F, 0x3B, 0x63)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF3, 0xF9)
RED = RGBColor(0xB3, 0x1B, 0x1B)
GREEN = RGBColor(0x1E, 0x6F, 0x3C)
AMBER = RGBColor(0x9C, 0x63, 0x00)
MORPH = RGBColor(0xC0, 0x50, 0x4D)
WATER = RGBColor(0x5B, 0x8F, 0xB0)
NONE = RGBColor(0x8A, 0x8A, 0x8A)
ABST = RGBColor(0xC5, 0xB8, 0xD9)
PALE_RED = RGBColor(0xFD, 0xE9, 0xE9)
PALE_GREEN = RGBColor(0xD8, 0xEC, 0xD8)
PALE_AMBER = RGBColor(0xFF, 0xF4, 0xD6)
BAND = RGBColor(0xD9, 0xE4, 0xF0)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.5)


def deck() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _run(shape, size, color, bold=False, align=PP_ALIGN.LEFT, line=1.0):
    tf = shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        para.line_spacing = line
        for r in para.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"


def box(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    return sh


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    return sh


def txt(slide, l, t, w, h, text, size, color=NAVY, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    sh = slide.shapes.add_textbox(l, t, w, h)
    sh.text_frame.word_wrap = True
    sh.text_frame.auto_size = None
    sh.text_frame.vertical_anchor = anchor
    sh.text_frame.text = text
    _run(sh, size, color, bold=bold, align=align, line=1.02)
    return sh


def header(slide, title, message):
    rect(slide, 0, 0, W, Inches(1.15), NAVY)
    txt(slide, M, Inches(0.12), W - 2 * M, Inches(0.5),
        title, 26, WHITE, bold=True)
    txt(slide, M, Inches(0.62), W - 2 * M, Inches(0.42),
        message, 16, RGBColor(0xC9, 0xDA, 0xEC), bold=False)


def bar(slide, l, t, w, h, fill, label, fg=WHITE, size=11):
    sh = rect(slide, l, t, w, h, fill, line=WHITE)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.text = label
    _run(sh, size, fg, bold=True, align=PP_ALIGN.CENTER, line=0.9)
    return sh


def pin(slide, x, bar_top, label, color=RED):
    stem = rect(slide, x - Emu(5000), bar_top - Inches(0.18), Emu(10000), Inches(0.18), color)
    txt(slide, x - Inches(0.9), bar_top - Inches(0.52), Inches(1.8), Inches(0.34),
        label, 11, color, bold=True, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------


def s01(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    txt(s, M, Inches(1.6), W - 2 * M, Inches(1.3),
        "Chemogenetic test of the Post ensemble",
        36, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, M, Inches(3.0), W - 2 * M, Inches(0.5),
        "One decision for my PI: which question do we test this year?",
        20, RGBColor(0xC9, 0xDA, 0xEC), align=PP_ALIGN.CENTER)
    locked = (
        "Already locked   ·   TRAP2 × Ai14   ·   tag at Post day 12   ·   "
        "wait 14 days for hM4Di   ·   DCZ   ·   Active and Passive   ·   ORBm first"
    )
    txt(s, M, Inches(5.8), W - 2 * M, Inches(0.5),
        locked, 14, RGBColor(0xA8, 0xC0, 0xD8), align=PP_ALIGN.CENTER)


def s02(prs):
    s = blank(prs)
    header(s, "The constraint",
           "Tag at day 12  +  14-day wait  =  first possible DCZ is day 26. Post is already over.")

    y = Inches(2.05)
    h = Inches(0.7)
    left, width = M, W - 2 * M
    # original 18-day protocol, equal-ish visual weights by phase not day
    segs = [
        ("Pre  water", 3, WATER),
        ("During  morphine", 5, MORPH),
        ("Post  morphine", 3, MORPH),
        ("Withdrawal  water", 3, WATER),
        ("Re-exp  morphine", 2, MORPH),
    ]
    days = [3, 5, 3, 3, 2]
    total = sum(days)
    x = left
    for (lab, d, fill), d in zip(segs, days):
        w = Emu(int(width * d / total))
        bar(s, x, y, w, h, fill, lab, size=12)
        x += w

    # day markers
    def frac(day):
        return (day - 0.5) / 18

    pin(s, left + Emu(int(width * frac(12))), y, "4-OHT  day 12")
    # DCZ wanted in Post - fail
    post_mid = left + Emu(int(width * (3 + 5 + 1.5) / 18))
    txt(s, post_mid - Inches(1.1), y + h + Inches(0.08), Inches(2.2), Inches(0.35),
        "DCZ here?  NO", 14, RED, bold=True, align=PP_ALIGN.CENTER)

    ready_x = left + width + Inches(0.02)
    # show day 26 off the bar
    txt(s, left, Inches(3.5), width, Inches(0.4),
        "day 1                                         12                                              18",
        12, GREY, align=PP_ALIGN.LEFT)

    box(s, M, Inches(4.3), W - 2 * M, Inches(2.3), PALE_RED, line=RED)
    txt(s, M + Inches(0.35), Inches(4.55), W - 2 * M - Inches(0.7), Inches(1.9),
        "hM4Di ready  =  day 26\n"
        "Post window  =  day 11–13\n"
        "Original sketch (DCZ on Post day 7) tests a receptor that is not there.",
        20, NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def s03(prs):
    s = blank(prs)
    header(s, "Two questions. This cohort can answer only one.",
           "14-day wait removes the left box. Choose the right box this year, or change the tool.")

    w = Inches(5.7)
    box(s, M, Inches(1.55), w, Inches(5.3), PALE_RED, line=RED)
    txt(s, M + Inches(0.25), Inches(1.75), w - Inches(0.5), Inches(0.7),
        "Q1   During Post", 22, RED, bold=True)
    txt(s, M + Inches(0.25), Inches(2.5), w - Inches(0.5), Inches(2.4),
        "Silence the tagged cells\nwhile the mouse is still taking morphine.\nDoes Post seeking fall?",
        18, NAVY, align=PP_ALIGN.LEFT)
    txt(s, M + Inches(0.25), Inches(5.15), w - Inches(0.5), Inches(1.3),
        "Needs receptor inside a 3-day window.\nAAV + 14-day wait cannot do this.",
        16, GREY)

    box(s, Inches(7.1), Inches(1.55), w, Inches(5.3), PALE_GREEN, line=GREEN)
    txt(s, Inches(7.35), Inches(1.75), w - Inches(0.5), Inches(0.7),
        "Q2   After abstinence", 22, GREEN, bold=True)
    txt(s, Inches(7.35), Inches(2.5), w - Inches(0.5), Inches(2.4),
        "Silence the same cells later.\nDoes the mouse still seek\nwhen put back in the box?",
        18, NAVY)
    txt(s, Inches(7.35), Inches(5.15), w - Inches(0.5), Inches(1.3),
        "Receptor is ready. Behaviour protocol\nmust be extended past day 18.",
        16, GREY)


def s04(prs):
    s = blank(prs)
    header(s, "Three different sessions. “No morphine” is not one condition.",
           "Pick the day-33 task. This changes what the experiment is allowed to claim.")

    cols = [
        ("Your Withdrawal\nday 14–16", WATER,
         "Pump ON\nWater comes out",
         "Water motivation\nafter morphine",
         "Not drug seeking"),
        ("Your Re-exposure\nday 17–18", MORPH,
         "Pump ON\nMorphine comes out",
         "Taking the drug again",
         "Not seeking"),
        ("Relapse test\n(not in your 18-day protocol)", NONE,
         "Pump OFF\nNothing comes out",
         "Cue / box seeking",
         "Standard relapse assay"),
    ]
    x = M
    w = Inches(3.9)
    gap = Inches(0.27)
    for title, col, mid, claim, note in cols:
        box(s, x, Inches(1.55), w, Inches(5.3), LIGHT, line=col)
        head = rect(s, x, Inches(1.55), w, Inches(1.15), col)
        head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        head.text_frame.word_wrap = True
        head.text_frame.text = title
        _run(head, 16, WHITE, bold=True, align=PP_ALIGN.CENTER, line=1.05)
        txt(s, x + Inches(0.2), Inches(2.9), w - Inches(0.4), Inches(1.3),
            mid, 18, NAVY, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), Inches(4.3), w - Inches(0.4), Inches(1.1),
            claim, 16, NAVY, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), Inches(5.55), w - Inches(0.4), Inches(0.9),
            note, 14, GREY, align=PP_ALIGN.CENTER)
        x += w + gap


def s05(prs):
    s = blank(prs)
    header(s, "All options",
           "A and B test Q2. C tries to keep Q1 by stretching morphine. D is not valid. E is next year.")

    rows = [
        ["", "What happens", "Pro", "Con"],
        ["A  Relapse test\nday 33",
         "4-OHT day 12\nDCZ once, day 33\nNothing delivered",
         "Receptor ready\nStandard relapse claim",
         "New session type\nNot your Withdrawal"],
        ["B  Your Withdrawal\nat day 33",
         "Same as A, but\nwater PR like day 14–16",
         "Matches your protocol\nComparable to TRAP data",
         "Water is available\nWeak as “seeking”"],
        ["C  Stretch Post",
         "Keep morphine until\n~day 26, then DCZ",
         "Tests Q1\nSilence while taking",
         "Drug history ≠ screen\nLong extra morphine"],
        ["D  Same mouse\nday 19 and 33",
         "DCZ twice in one mouse",
         "Fewer mice",
         "Day 19 = only +7 days\n2nd test is extinction"],
        ["E  New mouse line\n(Phase 2)",
         "TRAP2 × R26-hM4Di\n+ local DCZ cannula",
         "Ready in 48 h\nTrue 3-day Post test",
         "Breed 4–6 months\nNo Ai14 on same mouse"],
    ]
    nr, nc = len(rows), 4
    left, top, width = M, Inches(1.45), W - 2 * M
    row_h, hdr_h = Inches(0.88), Inches(0.36)
    gt = s.shapes.add_table(nr, nc, left, top, width, hdr_h + row_h * (nr - 1)).table
    col_w = [2.2, 3.3, 3.4, 3.4]
    tot = sum(col_w)
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Emu(int(width * cw / tot))
    gt.rows[0].height = hdr_h
    fills = [None, PALE_GREEN, PALE_AMBER, LIGHT, PALE_RED, LIGHT]
    for i, row in enumerate(rows):
        if i:
            gt.rows[i].height = row_h
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.text = val
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.04)
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = NAVY
            else:
                c.fill.fore_color.rgb = fills[i]
            for para in c.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT if j else PP_ALIGN.CENTER
                para.line_spacing = 0.95
                for r in para.runs:
                    r.font.size = Pt(12 if i else 13)
                    r.font.bold = i == 0 or j == 0
                    r.font.color.rgb = WHITE if i == 0 else NAVY
                    r.font.name = "Calibri"


def s06(prs):
    s = blank(prs)
    header(s, "Option A  —  the day-33 relapse test",
           "One DCZ injection. One test. Nothing comes out of the pump.")

    y = Inches(2.35)
    h = Inches(0.85)
    left, width = M, W - 2 * M
    segs = [
        ("1–10  Pre/During", 10, WATER),
        ("11–13  Post", 3, MORPH),
        ("14–18  Wd / Re-exp", 5, WATER),
        ("19–32  home cage", 14, ABST),
        ("33  TEST", 2, GREEN),
    ]
    total = sum(x[1] for x in segs)
    x = left
    for lab, d, fill in segs:
        w = Emu(int(width * d / total))
        fg = WHITE if fill in (MORPH, GREEN, NAVY) else NAVY
        bar(s, x, y, w, h, fill, lab, fg=fg, size=11)
        x += w

    t_x = left + Emu(int(width * (10 + 1.5) / total))
    pin(s, t_x, y, "4-OHT")
    dcz_x = left + Emu(int(width * (10 + 3 + 5 + 14 + 1) / total))
    pin(s, dcz_x, y, "DCZ + test")

    txt(s, M, Inches(3.6), W - 2 * M, Inches(0.4),
        "day 12                                              day 26 receptor ready                         day 33",
        13, GREY, align=PP_ALIGN.CENTER)

    box(s, M, Inches(4.3), Inches(6.0), Inches(2.3), PALE_GREEN, line=GREEN)
    txt(s, M + Inches(0.25), Inches(4.5), Inches(5.5), Inches(1.9),
        "PRO\nReceptor ready  (+21 days)\nClean relapse claim\nMatches OFC literature day",
        16, NAVY, bold=False, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(7.3), Inches(4.3), Inches(5.5), Inches(2.3), PALE_AMBER, line=AMBER)
    txt(s, Inches(7.55), Inches(4.5), Inches(5.05), Inches(1.9),
        "CON\nNot your Withdrawal task\nCannot say “incubation”\nDoes not test Q1",
        16, NAVY, anchor=MSO_ANCHOR.MIDDLE)


def s07(prs):
    s = blank(prs)
    header(s, "Option B  —  same calendar, your Withdrawal task",
           "Day 33 is water PR, like day 14–16. Pump is ON. Water comes out.")

    # reuse A timeline colours but last bar is water
    y = Inches(2.35)
    h = Inches(0.85)
    left, width = M, W - 2 * M
    segs = [
        ("1–10  Pre/During", 10, WATER),
        ("11–13  Post", 3, MORPH),
        ("14–18  Wd / Re-exp", 5, WATER),
        ("19–32  home cage", 14, ABST),
        ("33  water PR", 2, WATER),
    ]
    total = sum(x[1] for x in segs)
    x = left
    for lab, d, fill in segs:
        w = Emu(int(width * d / total))
        fg = WHITE if fill in (MORPH, GREEN, NAVY) else NAVY
        bar(s, x, y, w, h, fill, lab, fg=fg, size=11)
        x += w
    pin(s, left + Emu(int(width * 11.5 / total)), y, "4-OHT")
    pin(s, left + Emu(int(width * 33 / 34.0)), y, "DCZ + water")

    box(s, M, Inches(4.3), Inches(6.0), Inches(2.3), PALE_GREEN, line=GREEN)
    txt(s, M + Inches(0.25), Inches(4.5), Inches(5.5), Inches(1.9),
        "PRO\nSame task you already run\nLinks to TRAP Withdrawal data",
        16, NAVY, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(7.3), Inches(4.3), Inches(5.5), Inches(2.3), PALE_AMBER, line=AMBER)
    txt(s, Inches(7.55), Inches(4.5), Inches(5.05), Inches(1.9),
        "CON\nMouse can earn water\nHard to call this drug seeking\nReviewers will notice",
        16, NAVY, anchor=MSO_ANCHOR.MIDDLE)


def s08(prs):
    s = blank(prs)
    header(s, "Option C  —  stretch Post until the receptor is ready",
           "Morphine continues to ~day 26. Then DCZ. This is Q1, not relapse.")

    y = Inches(2.35)
    h = Inches(0.85)
    left, width = M, W - 2 * M
    segs = [
        ("1–10", 10, WATER),
        ("11–13  Post", 3, MORPH),
        ("14–25  still morphine", 12, MORPH),
        ("26  DCZ", 2, GREEN),
    ]
    total = sum(x[1] for x in segs)
    x = left
    for lab, d, fill in segs:
        w = Emu(int(width * d / total))
        fg = WHITE if fill in (MORPH, GREEN, NAVY) else NAVY
        bar(s, x, y, w, h, fill, lab, fg=fg, size=12)
        x += w
    pin(s, left + Emu(int(width * 11.5 / total)), y, "4-OHT")
    pin(s, left + Emu(int(width * 25.5 / total)), y, "DCZ on morphine")

    box(s, M, Inches(4.3), Inches(6.0), Inches(2.3), PALE_GREEN, line=GREEN)
    txt(s, M + Inches(0.25), Inches(4.5), Inches(5.5), Inches(1.9),
        "PRO\nTests the original hypothesis\nSilence while taking",
        16, NAVY, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(7.3), Inches(4.3), Inches(5.5), Inches(2.3), PALE_RED, line=RED)
    txt(s, Inches(7.55), Inches(4.5), Inches(5.05), Inches(1.9),
        "CON\n~2 extra weeks of morphine\nNot the protocol behind the screen\nCannot compare to current TRAP",
        16, NAVY, anchor=MSO_ANCHOR.MIDDLE)


def s09(prs):
    s = blank(prs)
    header(s, "Option D  —  do not use the same mouse twice",
           "Day 19 is too early for hM4Di. Day 33 is then an extinction retest.")

    # two mice vs one mouse visual
    box(s, M, Inches(1.55), Inches(6.0), Inches(5.3), PALE_RED, line=RED)
    txt(s, M + Inches(0.3), Inches(1.8), Inches(5.4), Inches(0.5),
        "Same mouse, two tests", 20, RED, bold=True)
    txt(s, M + Inches(0.3), Inches(2.5), Inches(5.4), Inches(3.8),
        "Day 19 DCZ  =  only 7 days after 4-OHT\n→ receptor not ready\n\n"
        "Day 33 test  =  second extinction\n→ seeking is already trained down\n\n"
        "If DCZ both days  →  no vehicle control",
        16, NAVY)

    box(s, Inches(7.3), Inches(1.55), Inches(5.5), Inches(5.3), PALE_GREEN, line=GREEN)
    txt(s, Inches(7.55), Inches(1.8), Inches(5.05), Inches(0.5),
        "Rule", 20, GREEN, bold=True)
    txt(s, Inches(7.55), Inches(2.5), Inches(5.05), Inches(3.8),
        "One mouse  =  one seeking test\n\n"
        "DCZ or vehicle, not both\n\n"
        "That is why A/B/C use\nseparate groups, not a crossover",
        16, NAVY)


def s10(prs):
    s = blank(prs)
    header(s, "Option E  —  next year, if Q1 is still the goal",
           "Transgenic hM4Di is on in 48 h. Then the original 3-day Post can be tested.")

    y = Inches(2.2)
    h = Inches(0.85)
    left, width = M, W - 2 * M
    segs = [
        ("1–10  Pre/During", 10, WATER),
        ("11–13  Post", 3, MORPH),
        ("12  4-OHT", 0.01, MORPH),
    ]
    bar(s, left, y, Emu(int(width * 10 / 13)), h, WATER, "1–10  Pre / During", fg=NAVY, size=14)
    bar(s, left + Emu(int(width * 10 / 13)), y, Emu(int(width * 3 / 13)), h, MORPH,
        "11–13  Post   ·   DCZ possible on day 13", fg=WHITE, size=14)
    pin(s, left + Emu(int(width * 11.5 / 13)), y, "4-OHT day 12")

    box(s, M, Inches(3.7), Inches(6.0), Inches(2.9), PALE_GREEN, line=GREEN)
    txt(s, M + Inches(0.25), Inches(3.9), Inches(5.5), Inches(2.5),
        "PRO\n48 h, not 14 days\nOriginal 3-day Post intact\nTogether with A: taking vs relapse",
        16, NAVY, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(7.3), Inches(3.7), Inches(5.5), Inches(2.9), PALE_AMBER, line=AMBER)
    txt(s, Inches(7.55), Inches(3.9), Inches(5.05), Inches(2.5),
        "CON\nStart breeding now\nRosa26: cannot also carry Ai14\nNeeds cannulae for region",
        16, NAVY, anchor=MSO_ANCHOR.MIDDLE)


def s11(prs):
    s = blank(prs)
    header(s, "Groups if you pick A or B",
           "Five groups. One test each. n = 12. Order 75 for ORBm. Same again later for BMAp.")

    rows = [
        ["#", "Mice", "Virus", "Day 33", "What it answers"],
        ["1", "Active", "hM4Di", "DCZ", "Does seeking fall?"],
        ["2", "Active", "hM4Di", "vehicle", "Same virus, no drug"],
        ["3", "Active", "mCherry", "DCZ", "Is it the ligand?"],
        ["4", "Passive", "hM4Di", "DCZ", "Does it fall without volition?"],
        ["5", "Passive", "hM4Di", "vehicle", "Passive baseline"],
    ]
    nr, nc = 6, 5
    left, top, width = M, Inches(1.5), W - 2 * M
    row_h, hdr_h = Inches(0.62), Inches(0.4)
    gt = s.shapes.add_table(nr, nc, left, top, width, hdr_h + row_h * 5).table
    col_w = [0.7, 1.8, 2.2, 2.0, 5.6]
    tot = sum(col_w)
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Emu(int(width * cw / tot))
    gt.rows[0].height = hdr_h
    for i, row in enumerate(rows):
        if i:
            gt.rows[i].height = row_h
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.text = val
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.08)
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY if i == 0 else (WHITE if i % 2 else LIGHT)
            for para in c.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER if j < 4 else PP_ALIGN.LEFT
                for r in para.runs:
                    r.font.size = Pt(16 if i else 14)
                    r.font.bold = i == 0 or j == 0
                    r.font.color.rgb = WHITE if i == 0 else NAVY
                    r.font.name = "Calibri"

    txt(s, M, Inches(6.35), W - 2 * M, Inches(0.55),
        "1 vs 2 = hypothesis     ·     1 vs 3 = ligand     ·     4 vs 5 = Active-only?",
        16, GREY, align=PP_ALIGN.CENTER)


def s12(prs):
    s = blank(prs)
    header(s, "Please pick",
           "Three ticks. That is the whole decision.")

    items = [
        ("1   Question",
         "Q2  after abstinence   (A or B, this year)\n"
         "Q1  during Post        (C this year, or E next year)"),
        ("2   Day-33 task  (only if Q2)",
         "A   nothing delivered   =  relapse / seeking\n"
         "B   water PR            =  your Withdrawal task"),
        ("3   Region order",
         "ORBm first, then BMAp\n"
         "75 mice ordered for ORBm"),
    ]
    y = Inches(1.5)
    for head, body in items:
        box(s, M, y, W - 2 * M, Inches(1.65), LIGHT, line=BAND)
        txt(s, M + Inches(0.3), y + Inches(0.12), W - 2 * M - Inches(0.6), Inches(0.4),
            head, 18, BLUE, bold=True)
        txt(s, M + Inches(0.3), y + Inches(0.52), W - 2 * M - Inches(0.6), Inches(1.0),
            body, 16, NAVY)
        y += Inches(1.8)

    txt(s, M, Inches(6.9), W - 2 * M, Inches(0.35),
        "My recommendation:  1 = Q2    2 = A    3 = ORBm first.  Start E breeding if Q1 still matters.",
        14, GREY, align=PP_ALIGN.CENTER)


def main() -> None:
    prs = deck()
    for fn in (s01, s02, s03, s04, s05, s06, s07, s08, s09, s10, s11, s12):
        fn(prs)
    out = Path(__file__).resolve().parents[2] / "outputs" / "Chemo_PI_decisions.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    # also replace the wordy deck path the PI was given last time
    alt = Path(__file__).resolve().parents[2] / "outputs" / "NextPhase_Plan_ORBm_BMAp_TRAP.pptx"
    prs.save(alt)
    print(f"[DONE] {out}")
    print(f"[DONE] {alt}  ({len(list(prs.slides))} slides)")


if __name__ == "__main__":
    main()
