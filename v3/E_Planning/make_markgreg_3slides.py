#!/usr/bin/env python
"""3-slide PI brief for Mark & Greg. One typeface, two sizes, two colors."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1F, 0x2A, 0x44)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)
LINE = RGBColor(0xD0, 0xD0, 0xD0)
FILL = RGBColor(0xF4, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# Single accent — used only on the three confirmation ticks
ASK = RGBColor(0x1F, 0x2A, 0x44)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.55)
TITLE_SIZE = 22
BODY = 14
SMALL = 12  # tables only


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def style(tf, size, color, bold=False, align=PP_ALIGN.LEFT, line=1.05):
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        para.line_spacing = line
        para.space_after = Pt(0)
        for r in para.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"


def rect(s, l, t, w, h, fill, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def tb(s, l, t, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP):
    sh = s.shapes.add_textbox(l, t, w, h)
    sh.text_frame.word_wrap = True
    sh.text_frame.vertical_anchor = anchor
    sh.text_frame.text = text
    style(sh.text_frame, size, color, bold=bold, align=align)
    return sh


def header(s, title):
    rect(s, 0, 0, W, Inches(0.85), NAVY)
    tb(s, M, Inches(0.22), W - 2 * M, Inches(0.5), title, TITLE_SIZE, WHITE, bold=True)


def table(s, rows, l, t, w, col_w, row_h=Inches(0.42), hdr_h=Inches(0.38)):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, l, t, w, hdr_h + row_h * (nr - 1)).table
    tot = sum(col_w)
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Emu(int(w * cw / tot))
    gt.rows[0].height = hdr_h
    for i, row in enumerate(rows):
        if i:
            gt.rows[i].height = row_h
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.text = str(val)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.04)
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY if i == 0 else (WHITE if i % 2 else FILL)
            for para in c.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT if j else PP_ALIGN.LEFT
                para.line_spacing = 1.0
                for r in para.runs:
                    r.font.size = Pt(SMALL)
                    r.font.bold = i == 0
                    r.font.color.rgb = WHITE if i == 0 else INK
                    r.font.name = "Calibri"
    return gt


def s01(prs):
    s = blank(prs)
    header(s, "Please confirm three things")
    tb(s, M, Inches(1.05), W - 2 * M, Inches(0.35),
       "HEAL  ·  ORBm and BMAp  ·  meeting brief  ·  Excel gene list attached",
       BODY, MUTED)

    items = [
        ("1", "Gene list for Xenium",
         "Approve the attached list.\n"
         "Minimum to order: 77 genes (priority 1).\n"
         "Full list: 168 genes (BMAp 124, ORBm 128, 84 shared)."),
        ("2", "Xenium experiment",
         "3 Active + 3 Passive. Tag at Post (day 13).\n"
         "Read out after rest, on a seeking test (pump off).\n"
         "ORBm and BMAp from the same mouse."),
        ("3", "Causality route  (pick one on slide 3)",
         "A  Silence the TRAPped Post ensemble (DREADD).\n"
         "B  Drug the GPCR that Xenium names (no TRAP wait).\n"
         "C  Silence the region, not the ensemble (fallback)."),
    ]
    x = M
    w = Inches(3.95)
    gap = Inches(0.22)
    for num, head, body in items:
        top = Inches(1.55)
        box = rect(s, x, top, w, Inches(5.2), WHITE, line=LINE)
        # number bar
        bar = rect(s, x, top, w, Inches(0.7), NAVY)
        bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bar.text_frame.text = f"  {num}   {head}"
        style(bar.text_frame, BODY, WHITE, bold=True)
        tb(s, x + Inches(0.2), top + Inches(0.95), w - Inches(0.4), Inches(3.4),
           body, BODY, INK)
        tb(s, x + Inches(0.2), top + Inches(4.5), w - Inches(0.4), Inches(0.45),
           "☐  yes          ☐  discuss", SMALL, MUTED)
        x += w + gap


def s02(prs):
    s = blank(prs)
    header(s, "Xenium + TRAP  —  what we would run")
    tb(s, M, Inches(1.05), W - 2 * M, Inches(0.35),
       "Same 18-day task you already know. Only addition: 4-OHT at Post, then rest, then one seeking test.",
       BODY, MUTED)

    # one-row timeline, two greys
    y = Inches(1.6)
    h = Inches(0.5)
    left = M
    width = W - 2 * M
    segs = [
        (10, "days 1–10  task", FILL, INK),
        (3, "Post  morphine", NAVY, WHITE),
        (5, "days 14–18", FILL, INK),
        (12, "rest  (protein)", FILL, INK),
        (2, "test + Xenium", NAVY, WHITE),
    ]
    total = sum(a[0] for a in segs)
    x = left
    for d, lab, fill, fg in segs:
        ww = Emu(int(width * d / total))
        b = rect(s, x, y, ww, h, fill, line=WHITE)
        b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        b.text_frame.word_wrap = True
        b.text_frame.text = lab
        style(b.text_frame, SMALL, fg, bold=True, align=PP_ALIGN.CENTER)
        x += ww
    tb(s, M, Inches(2.18), Inches(6), Inches(0.3),
       "4-OHT on Post day 3  (calendar day 13)", SMALL, MUTED)

    rows = [
        ["", "Specification"],
        ["Mice", "TRAP2 × Ai14   ·   3 Active, 3 Passive"],
        ["Tag", "4-OHT after the last Post session"],
        ["Readout", "Seeking test (pump off), then tissue. ORBm + BMAp."],
        ["Panel", "77 genes if we must cut. 168 if the add-on budget allows."],
        ["Why 77 vs 168", "77 still IDs all 20 cell types + TRAP tag + backbone. 168 adds extra GPCRs."],
        ["Attached", "FINAL_ordering_BMAp_ORBm_TRAP_HANSOL.xlsx"],
    ]
    table(s, rows, M, Inches(2.6), W - 2 * M, [2.4, 9.4],
          row_h=Inches(0.48), hdr_h=Inches(0.34))
    tb(s, M, Inches(6.85), W - 2 * M, Inches(0.35),
       "Confirm:  order this list  ·  n = 6 is enough for a cell-type map, not for a group test.",
       BODY, INK, bold=True)


def s03(prs):
    s = blank(prs)
    header(s, "Causality  —  pick one path")
    tb(s, M, Inches(1.02), W - 2 * M, Inches(0.32),
       "TRAP protein needs ~14 days. That is why A cannot silence the first Post. B and C do not wait on TRAP.",
       BODY, MUTED)

    rows = [
        ["", "A  TRAP + DREADD", "B  GPCR ligand", "C  Region DREADD"],
        ["What we turn off",
         "Cells tagged at Post",
         "One receptor, whole brain or local",
         "ORBm or BMAp cells, not the ensemble"],
        ["18-day task",
         "Unchanged, then rest, then test",
         "Unchanged. Drug on the test day",
         "Unchanged. No 4-OHT"],
        ["Needs Xenium first?",
         "No, but Xenium names the cell type",
         "Yes — Xenium must name the GPCR",
         "No"],
        ["Surgery", "Yes (AAV)", "No", "Yes (AAV)"],
        ["n / region", "16 (8 Active, 8 Passive)", "same 16, no virus", "16"],
        ["Pro",
         "Tests the tagged ensemble",
         "No wait, no surgery, a real drug",
         "Works even if TRAP tagging is weak"],
        ["Con",
         "Cannot test during the first Post",
         "Not ensemble-specific",
         "Not ensemble-specific"],
        ["Please tick", "☐", "☐", "☐"],
    ]
    table(s, rows, M, Inches(1.4), W - 2 * M, [2.1, 3.5, 3.5, 3.2],
          row_h=Inches(0.52), hdr_h=Inches(0.36))

    tb(s, M, Inches(6.7), W - 2 * M, Inches(0.5),
       "If A: recommended timeline is rest after day 18, then DCZ vs vehicle on the test days "
       "(not a second 18-day run; not extra morphine). Details in the long draft if needed.",
       SMALL, MUTED)


def main() -> None:
    prs = deck()
    s01(prs)
    s02(prs)
    s03(prs)
    out = Path(__file__).resolve().parents[2] / "outputs" / "HEAL_MarkGreg_3slides.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"[DONE] {out}")


if __name__ == "__main__":
    main()
